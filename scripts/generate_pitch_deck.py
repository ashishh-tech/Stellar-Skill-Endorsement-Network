from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color Palette
DARK_BG = RGBColor(13, 17, 23)
ACCENT_BLUE = RGBColor(88, 166, 255)
ACCENT_PURPLE = RGBColor(188, 140, 255)
TEXT_WHITE = RGBColor(240, 246, 252)
TEXT_MUTED = RGBColor(139, 148, 158)
CARD_BG = RGBColor(22, 27, 34)

def add_base_slide(title_text, subtitle_text=""):
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = DARK_BG
    bg.line.color.rgb = DARK_BG
    
    # Title Box
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.6), Inches(11.7), Inches(1.2))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = ACCENT_BLUE
    
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(14)
        p2.font.color.rgb = TEXT_MUTED
        
    return slide

# Slide 1: Cover
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
bg1.fill.solid()
bg1.fill.fore_color.rgb = DARK_BG
bg1.line.color.rgb = DARK_BG

tbox = slide1.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.5))
tf1 = tbox.text_frame
tf1.word_wrap = True
p = tf1.paragraphs[0]
p.text = "🌟 Skill Endorsement Network"
p.font.size = Pt(40)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE

p2 = tf1.add_paragraph()
p2.text = "A Sybil-Resistant, On-Chain Reputation Graph on Stellar Soroban"
p2.font.size = Pt(22)
p2.font.color.rgb = ACCENT_PURPLE

p3 = tf1.add_paragraph()
p3.text = "\nLevel 7 — The Founder Belt Startup Pitch Deck\nLive Mainnet: https://skill-endorsement-network.netlify.app/\nFounder: Ashish (@ashishh-tech)"
p3.font.size = Pt(14)
p3.font.color.rgb = TEXT_MUTED

# Slide 2: Problem Statement
slide2 = add_base_slide("📌 The Problem: Broken Web2 & Web3 Reputation", "Traditional endorsement systems lack Sybil resistance and mathematical weight")
c1 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(3.6), Inches(4.5))
c1.fill.solid()
c1.fill.fore_color.rgb = CARD_BG
c1.line.color.rgb = ACCENT_PURPLE
tf = c1.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Zero Sybil Resistance\n"
p.font.bold = True
p.font.size = Pt(18)
p.font.color.rgb = ACCENT_BLUE
p_body = tf.add_paragraph()
p_body.text = "• Anyone can create unlimited bot accounts.\n• Self-endorsements go unchecked on Web2 platforms like LinkedIn.\n• Zero cryptographic proof of actual peer validation."
p_body.font.size = Pt(13)
p_body.font.color.rgb = TEXT_WHITE

c2 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.8), Inches(2.0), Inches(3.6), Inches(4.5))
c2.fill.solid()
c2.fill.fore_color.rgb = CARD_BG
c2.line.color.rgb = ACCENT_PURPLE
tf = c2.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Unweighted Flat Scores\n"
p.font.bold = True
p.font.size = Pt(18)
p.font.color.rgb = ACCENT_BLUE
p_body = tf.add_paragraph()
p_body.text = "• Senior architect endorsement carries the same weight as a brand-new account.\n• No dynamic reputation scoring or trust propagation.\n• Subject to reciprocal vote trading."
p_body.font.size = Pt(13)
p_body.font.color.rgb = TEXT_WHITE

c3 = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.8), Inches(2.0), Inches(3.6), Inches(4.5))
c3.fill.solid()
c3.fill.fore_color.rgb = CARD_BG
c3.line.color.rgb = ACCENT_PURPLE
tf = c3.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Centralized Data Silos\n"
p.font.bold = True
p.font.size = Pt(18)
p.font.color.rgb = ACCENT_BLUE
p_body = tf.add_paragraph()
p_body.text = "• User reputation is trapped inside proprietary platforms.\n• Cannot be verified by external smart contracts, DAOs, or recruiters.\n• No immutable audit trail on-chain."
p_body.font.size = Pt(13)
p_body.font.color.rgb = TEXT_WHITE

# Slide 3: The Solution
slide3 = add_base_slide("💡 The Solution: Soroban Reputation Graph", "Trust-weighted skill endorsements executed via atomic inter-contract calls on Stellar")
box3 = slide3.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf3 = box3.text_frame
tf3.word_wrap = True
points = [
    ("⚡ Inter-Contract Reputation Weighting:", " Endorsements dynamically query the endorser's real-time score: Weight = max(Reputation / 10, 1)."),
    ("🛡️ Protocol-Enforced Anti-Sybil Guards:", " Smart contract blocks self-endorsements and duplicate endorsements at the execution layer."),
    ("🌐 Multi-Wallet Accessibility:", " Seamless support for Freighter, Albedo, Hana, and xBull via StellarWalletsKit."),
    ("📊 Transparent Telemetry & Event Streaming:", " Real-time RPC latency monitoring and immutable ledger event feeds.")
]
for title, desc in points:
    p = tf3.add_paragraph()
    p.text = title + desc
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(14)

# Slide 4: Market Opportunity & Traction
slide4 = add_base_slide("📈 Traction & Market Opportunity", "Proven demand across 55+ verified mainnet users and 140+ on-chain transactions")
m1 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(2.7), Inches(4.5))
m1.fill.solid()
m1.fill.fore_color.rgb = CARD_BG
m1.line.color.rgb = ACCENT_BLUE
tf = m1.text_frame
p = tf.paragraphs[0]
p.text = "55+ Users\n\n"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = ACCENT_BLUE
p_body = tf.add_paragraph()
p_body.text = "Verified mainnet users across 24 countries with distinct on-chain profiles."
p_body.font.size = Pt(13)
p_body.font.color.rgb = TEXT_WHITE

m2 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.8), Inches(2.0), Inches(2.7), Inches(4.5))
m2.fill.solid()
m2.fill.fore_color.rgb = CARD_BG
m2.line.color.rgb = ACCENT_PURPLE
tf = m2.text_frame
p = tf.paragraphs[0]
p.text = "140+ Tx\n\n"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = ACCENT_PURPLE
p_body = tf.add_paragraph()
p_body.text = "On-chain Soroban transactions with 3.2s confirmation and <0.00004 XLM fees."
p_body.font.size = Pt(13)
p_body.font.color.rgb = TEXT_WHITE

m3 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(2.0), Inches(2.7), Inches(4.5))
m3.fill.solid()
m3.fill.fore_color.rgb = CARD_BG
m3.line.color.rgb = RGBColor(63, 185, 80)
tf = m3.text_frame
p = tf.paragraphs[0]
p.text = "4.94 / 5.0\n\n"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(63, 185, 80)
p_body = tf.add_paragraph()
p_body.text = "Customer satisfaction score across 55 detailed user feedback responses."
p_body.font.size = Pt(13)
p_body.font.color.rgb = TEXT_WHITE

m4 = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9.8), Inches(2.0), Inches(2.7), Inches(4.5))
m4.fill.solid()
m4.fill.fore_color.rgb = CARD_BG
m4.line.color.rgb = RGBColor(210, 153, 34)
tf = m4.text_frame
p = tf.paragraphs[0]
p.text = "68+ Social\n\n"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(210, 153, 34)
p_body = tf.add_paragraph()
p_body.text = "Organic X/Twitter audience with 14.2K post impressions and high viral engagement."
p_body.font.size = Pt(13)
p_body.font.color.rgb = TEXT_WHITE

# Slide 5: Business Model & Monetization
slide5 = add_base_slide("💼 Startup Business Model & Monetization on Stellar", "Three revenue streams driving sustainable long-term venture growth")
bm_box = slide5.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf = bm_box.text_frame
tf.word_wrap = True
streams = [
    ("1. B2B Enterprise Talent Verification API:", " Monthly subscription for Web3 talent recruiters and venture funds to query verified developer competence ($99 - $499/mo)."),
    ("2. Organization Attestation Fees:", " Developer bootcamps and Stellar projects pay micro-fees in XLM to issue verified alumni credential badges."),
    ("3. Soroban Micro-Bounties & Escrows:", " Peer code reviewers receive reputation stakes and XLM micro-payouts upon successful task completion.")
]
for title, desc in streams:
    p = tf.add_paragraph()
    p.text = title + desc
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(16)

# Slide 6: Roadmap & Founder Vision
slide6 = add_base_slide("🚀 30 / 60 / 90-Day Product Roadmap", "Scaling from hackathon project to enterprise-grade Stellar startup")
rm_box = slide6.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(11.7), Inches(5.0))
tf = rm_box.text_frame
tf.word_wrap = True
roadmap = [
    ("📍 Phase 1 (Completed):", " Stellar Soroban Mainnet Launch • 55+ Verified Users • Multi-Wallet • Sybil Simulator • 20/20 Test Suite"),
    ("📍 Phase 2 (Month 2 - In Progress):", " Mobile PWA with Biometric Passkeys • Batch Endorsement Smart Contract (`endorse_batch`) • Organization Portal"),
    ("📍 Phase 3 (Month 3 - Upcoming):", " W3C Verifiable Credentials / DID Resolver • Cross-Chain Reputation Bridges • Enterprise Recruiter Dashboard")
]
for title, desc in roadmap:
    p = tf.add_paragraph()
    p.text = title + desc
    p.font.size = Pt(16)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(16)

prs.save("Skill_Endorsement_Network_Pitch_Deck.pptx")
print("Successfully generated Skill_Endorsement_Network_Pitch_Deck.pptx")
